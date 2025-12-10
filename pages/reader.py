import streamlit as st
import app.pages as pages
import app.utils as utils
import app.prompts as prompts
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO


# --- UI helper: centered "OR" header with lines ---
def or_header(text: str):
    st.markdown(
        """
        <style>
        .or-header{display:flex;align-items:center;gap:.75rem;margin:.5rem 0 1rem;}
        .or-header:before,.or-header:after{content:"";flex:1;border-top:1px solid rgba(128,128,128,.35);}
        .or-header .or-text{font-weight:600;opacity:.85;}
        </style>
        """,
        unsafe_allow_html=True,
    )
    st.markdown(f"<div class='or-header'><span class='or-text'>{text}</span></div>", unsafe_allow_html=True)

# App title
pages.show_home()
pages.show_sidebar()

st.header("🔍Style Reader")

# st.session_state.exampleText = st.text_area(
#     ":blue[**Reference Input for BSP Writing Style :**]", st.session_state.exampleText, 200
# )

# uploaded_files = st.file_uploader(
#     ":blue[**Upload Example Files:**]", 
#     type=["pdf", "docx", "pptx"], 
#     accept_multiple_files=True,
#     help="Upload PDF, Word, or PowerPoint files"
# )

or_header("Input the Contents or Upload the File for BSP Style Reading")

# --- two-column layout (Col 1 / Col 2) ---
col1, col2 = st.columns([3, 2], gap="small")

with col1:
    
    st.session_state.content = st.text_area(
        ":blue[**Input Content:**]",
        st.session_state.content,
        height=160,
        key="content_input",
    )

with col2:
    
    uploaded_files = st.file_uploader(
        ":blue[**Upload Files:**]",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
        help="Upload PDF, Word, or PowerPoint files (Max 50MB per file recommended)",
        key="content_upload",
    )
    
    # Show file size warning
    if uploaded_files:
        total_size = sum(f.size for f in uploaded_files) / (1024 * 1024)  # Convert to MB
        if total_size > 100:
            st.warning(f"⚠️ Total file size: {total_size:.1f}MB. Large files may take longer to process.")

# Additional Prompt Instruction (Optional)
st.session_state.setdefault("additional_instruction_reader", "")
st.session_state.additional_instruction_reader = st.text_area(
    ":blue[**Additional Prompt Instruction (Optional):**]",
    st.session_state.additional_instruction_reader,
    height=100,
    key="additional_instruction_reader_input",
    help="Add any additional instructions to combine with the extracted style...",
)
st.caption("This will be appended to the extracted style when saved.")

# Extract text from uploaded files
extracted_text = ""
if uploaded_files:
    progress_bar = st.progress(0, text="Processing uploaded files...")
    
    for idx, uploaded_file in enumerate(uploaded_files):
        try:
            file_type = uploaded_file.name.split('.')[-1].lower()
            file_size_mb = uploaded_file.size / (1024 * 1024)
            
            # Update progress
            progress_bar.progress(
                (idx) / len(uploaded_files),
                text=f"Processing {uploaded_file.name} ({file_size_mb:.1f}MB)..."
            )
            
            # Read file content once
            file_content = uploaded_file.read()
            
            if not file_content:
                st.warning(f"⚠️ {uploaded_file.name} is empty or couldn't be read. Skipping.")
                continue
            
            if file_type == 'pdf':
                pdf_reader = PdfReader(BytesIO(file_content))
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        extracted_text += text + "\n"
                    
            elif file_type == 'docx':
                doc = Document(BytesIO(file_content))
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        extracted_text += paragraph.text + "\n"
                        
            elif file_type == 'pptx':
                prs = Presentation(BytesIO(file_content))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            extracted_text += shape.text + "\n"
            
            # Clear file content from memory
            del file_content
            
        except Exception as e:
            st.error(f"❌ Error processing {uploaded_file.name}: {str(e)}")
            continue
    
    # Complete progress
    progress_bar.progress(1.0, text="✓ All files processed!")
    
    # Clear progress bar after a moment
    import time
    time.sleep(0.5)
    progress_bar.empty()


# ----------------------------
# PICK EXACTLY ONE CONTENT SOURCE + PREVIEW
# ----------------------------
has_input = bool(st.session_state.content.strip())
has_uploads = bool(extracted_text.strip())

# Let the user choose if both are present; otherwise auto-pick
if has_input and has_uploads:
    source = st.radio(
        ":blue[**Choose content source**]",
        ["Uploaded files", "Manual input"],
        horizontal=True,
        index=0,
        help="Use either the text you typed OR the text extracted from the uploaded files."
    )
elif has_uploads:
    source = "Uploaded files"
elif has_input:
    source = "Manual input"
else:
    source = None

# Decide combined_text and show a preview when using uploads
if source == "Uploaded files":
    combined_text = extracted_text  # keep full unicode; your PDF builder handles fonts
    with st.expander("📄 Preview: Extracted text from uploaded files", expanded=True):
        st.text_area(
            "Extracted Text",
            combined_text,
            height=240,
            key="uploaded_text_preview",
        )
elif source == "Manual input":
    combined_text = st.session_state.content
else:
    combined_text = ""
    st.markdown(
    """
    <div class="bsp-alert-red" role="alert">
      <strong>Heads up:</strong> Provide content — either type in the left box or upload a file on the right.
    </div>
    <style>
      .bsp-alert-red{
        padding:12px 14px;
        margin: 4px 0 10px;
        border-radius:10px;
        border:1px solid rgba(220,53,69,.35);
        background: rgba(220,53,69,.08); /* light red */
        font-size: 0.95rem;
      }
      .bsp-alert-red strong{
        color:#b02a37; /* dark red for emphasis */
      }
    </style>
    """,
    unsafe_allow_html=True,
)

st.session_state.styleName = st.text_input(
    ":blue[**Input Style Name:**]", st.session_state.styleName, 100
)


# Combine text area and extracted content
#combined_text = st.session_state.exampleText + "\n" + extracted_text.encode("ascii", errors="ignore").decode("ascii")

if st.button(
    ":blue[**Extract Writing Style**]",
    key="extract",
   # disabled=combined_text.strip() == "" or st.session_state.styleName == "",
    disabled=(
    combined_text.strip() == ""
    or st.session_state.styleName == ""
    #or st.session_state.content == ""
)
):
    with st.spinner("Processing..."):
        # Check if style name already exists
        if utils.check_style(st.session_state.styleName):
            st.session_state["extraction_error"] = f"Style name '{st.session_state.styleName}' already exists. Please choose a different name."
            st.session_state["extraction_success"] = False
        else:
            style = prompts.extract_style(combined_text, False)
            utils.save_style(style, combined_text)
            st.session_state["extracted_style"] = style
            st.session_state["extracted_style_name"] = st.session_state.styleName
            st.session_state["extraction_success"] = True
            st.session_state["extraction_error"] = None

# Display extraction results persistently
if st.session_state.get("extraction_success"):
    with st.container(border=True):
        st.success(f"✅ Style '{st.session_state.get('extracted_style_name')}' extracted and saved successfully!")
        
        with st.expander("📋 View Extracted Style", expanded=False):
            st.text_area(
                "Extracted Writing Style",
                st.session_state.get("extracted_style", ""),
                height=300,
                key="extracted_style_display",
            )

if st.session_state.get("extraction_error"):
    with st.container(border=True):
        st.error(st.session_state["extraction_error"])
