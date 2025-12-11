import streamlit as st
import app.pages as pages
import app.utils as utils
import app.prompts as prompts
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from io import BytesIO

# --- NEW imports ---
from datetime import datetime
from io import BytesIO

# PDF
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
import os

# --- UI helper: centered "OR" header with lines ---
def or_header(text: str):
    st.markdown(
        """
        <style>
        .or-header{display:flex;align-items:center;gap:.75rem;margin:.5rem 0 1rem;}
        .or-header:before,.or-header:after{content:"";flex:1;border-top:1px solid rgba(128,128,128,.35);}
        .or-header .or-text{font-weight:600;opacity:.85;}
        </style>
        """,
        unsafe_allow_html=True,
    )
    st.markdown(f"<div class='or-header'><span class='or-text'>{text}</span></div>", unsafe_allow_html=True)

# App title
pages.show_home()
pages.show_sidebar()

# Home page
st.header("📝Style Writer")

# # Content input
# st.session_state.content = st.text_area(
#     ":blue[**Content Input for BSP Writing Style:**]", st.session_state.content, 200
# )

# uploaded_files = st.file_uploader(
#     ":blue[**Upload Content Files:**]", 
#     type=["pdf", "docx", "pptx"], 
#     accept_multiple_files=True,
#     help="Upload PDF, Word, or PowerPoint files"
# )
or_header("Input the Contents or Upload the File for BSP Style Writing")

# --- two-column layout (Col 1 / Col 2) ---
col1, col2 = st.columns([3, 2], gap="small")

with col1:
    
    st.session_state.content = st.text_area(
        ":blue[**Input Content:**]",
        st.session_state.content,
        height=160,
        key="content_input",
    )

    MIN_LEN, MAX_LEN, DEFAULT_LEN = 20, 75_000, 1_000

    # One source of truth
    st.session_state.setdefault("max_len", DEFAULT_LEN)
    st.session_state.setdefault("last_updated", None)

    def _update_from_slider():
        st.session_state.last_updated = "slider"
        st.session_state.max_len = st.session_state.max_len_slider

    def _update_from_input():
        st.session_state.last_updated = "input"
        v = st.session_state.max_len_input
        # Coerce + clamp
        try:
            v = int(v)
        except Exception:
            v = MIN_LEN
        st.session_state.max_len = max(MIN_LEN, min(MAX_LEN, v))

    # Keep both widgets synced to the source of truth (avoid ping-pong)
    if st.session_state.last_updated != "slider":
        st.session_state["max_len_slider"] = st.session_state.max_len
    if st.session_state.last_updated != "input":
        st.session_state["max_len_input"] = st.session_state.max_len

    col_slider, col_input = st.columns([3, 1])

    with col_slider:
        st.slider(
            ":blue[**Output Maximum Character Length (75,000 Maximum):**]",
            min_value=MIN_LEN,
            max_value=MAX_LEN,
            key="max_len_slider",
            on_change=_update_from_slider,
            disabled=False,
        )

    with col_input:
        st.number_input(   # use number_input for clean integer UX
            "No. of Characters",
            min_value=MIN_LEN,
            max_value=MAX_LEN,
            step=1,
            key="max_len_input",
            on_change=_update_from_input,
        )

    # Use this in your app
    max_output_length = st.session_state.max_len

with col2:
    
    uploaded_files = st.file_uploader(
        ":blue[**Upload Files:**]",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
        help="Upload PDF, Word, or PowerPoint files (Max 50MB per file recommended)",
        key="content_upload",
    )
    
    # Show file size warning
    if uploaded_files:
        total_size = sum(f.size for f in uploaded_files) / (1024 * 1024)  # Convert to MB
        if total_size > 100:
            st.warning(f"⚠️ Total file size: {total_size:.1f}MB. Large files may take longer to process.")

# Additional Prompt Instruction (Optional)
st.session_state.setdefault("additional_instruction", "")
st.session_state.additional_instruction = st.text_area(
    ":blue[**Additional Prompt Instruction (Optional):**]",
    st.session_state.additional_instruction,
    height=100,
    key="additional_instruction_input",
    help="Add any specific instructions for rewriting...",
)
st.caption("This will be added to the prompt sent to the AI for rewriting.")

# Extract text from uploaded files
extracted_text = ""
if uploaded_files:
    progress_bar = st.progress(0, text="Processing uploaded files...")
    
    for idx, uploaded_file in enumerate(uploaded_files):
        try:
            file_type = uploaded_file.name.split('.')[-1].lower()
            file_size_mb = uploaded_file.size / (1024 * 1024)
            
            # Update progress
            progress_bar.progress(
                (idx) / len(uploaded_files),
                text=f"Processing {uploaded_file.name} ({file_size_mb:.1f}MB)..."
            )
            
            # Read file content once
            file_content = uploaded_file.read()
            
            if not file_content:
                st.warning(f"⚠️ {uploaded_file.name} is empty or couldn't be read. Skipping.")
                continue
            
            if file_type == 'pdf':
                pdf_reader = PdfReader(BytesIO(file_content))
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        extracted_text += text + "\n"
                    
            elif file_type == 'docx':
                doc = Document(BytesIO(file_content))
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        extracted_text += paragraph.text + "\n"
                        
            elif file_type == 'pptx':
                prs = Presentation(BytesIO(file_content))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            extracted_text += shape.text + "\n"
            
            # Clear file content from memory
            del file_content
            
        except Exception as e:
            st.error(f"❌ Error processing {uploaded_file.name}: {str(e)}")
            continue
    
    # Complete progress
    progress_bar.progress(1.0, text="✓ All files processed!")
    
    # Clear progress bar after a moment
    import time
    time.sleep(0.5)
    progress_bar.empty()


# ----------------------------
# PICK EXACTLY ONE CONTENT SOURCE + PREVIEW
# ----------------------------
has_input = bool(st.session_state.content.strip())
has_uploads = bool(extracted_text.strip())

# Let the user choose if both are present; otherwise auto-pick
if has_input and has_uploads:
    source = st.radio(
        ":blue[**Choose content source**]",
        ["Uploaded files", "Manual input"],
        horizontal=True,
        index=0,
        help="Use either the text you typed OR the text extracted from the uploaded files."
    )
elif has_uploads:
    source = "Uploaded files"
elif has_input:
    source = "Manual input"
else:
    source = None

# Decide content_all and show a preview when using uploads
if source == "Uploaded files":
    content_all = extracted_text  # keep full unicode; your PDF builder handles fonts
    with st.expander("📄 Preview: Extracted text from uploaded files", expanded=True):
        st.text_area(
            "Extracted Text",
            content_all,
            height=240,
            key="uploaded_text_preview",
        )
elif source == "Manual input":
    content_all = st.session_state.content
else:
    content_all = ""
    st.markdown(
    """
    <div class="bsp-alert-red" role="alert">
      <strong>Heads up:</strong> Provide content — either type in the left box or upload a file on the right.
    </div>
    <style>
      .bsp-alert-red{
        padding:12px 14px;
        margin: 4px 0 10px;
        border-radius:10px;
        border:1px solid rgba(220,53,69,.35);
        background: rgba(220,53,69,.08); /* light red */
        font-size: 0.95rem;
      }
      .bsp-alert-red strong{
        color:#b02a37; /* dark red for emphasis */
      }
    </style>
    """,
    unsafe_allow_html=True,
)

# Combine text area and extracted content
#content_all = st.session_state.content + "\n" + extracted_text.encode("ascii", errors="ignore").decode("ascii")

# Extracting the styles and creating combined display options
styles_data = utils.get_styles()
style_options = [item['name'] for item in styles_data]
selected_style = st.selectbox(":blue[**Select a Style:**]", options=style_options, index=None)

# Assigning the selected style to the session state
if selected_style:
    # Find the matching style data
    filtered = next(
        (item for item in styles_data if str(item["name"]) == selected_style), None
    )
    
    if filtered:
        st.session_state.style = filtered["style"]
        st.session_state.example = filtered["example"]
        st.session_state.styleId = selected_style
        
# st.session_state.style = st.text_area(":blue[**Style:**]", st.session_state.style)

# Show the example style
guidelines = st.session_state.locals.get("relevant_guidelines", {})
guidelines_summary = st.session_state.locals.get("guideline_summaries", {}) 
selected_guidelines = []

st.write(":blue[**Select Editorial Style Guides:**]")

# Tooltip for guideline summary in the UI
def render_guideline_checkbox(section_name: str, content: str, col_key_prefix: str):
    default_checked = section_name in ["ACRONYMS AND ABBREVIATIONS","CAPITALIZATION","NUMBERS","PUNCTUATION","SPECIAL CHARACTERS","COMMON GRAMMATICAL ERRORS", "LATIN ABBREVIATIONS","DOCUMENT SPECIFICATIONS","WRITING LETTERS",
    "Common acronyms and abbreviations"]
    tooltip = guidelines_summary.get(section_name, None)  # one-sentence summary for hover
    if st.checkbox(
        section_name,
        value=default_checked,
        key=f"{col_key_prefix}_{section_name}",
        help=tooltip  # <-- hover tooltip appears on the ⓘ icon and on hover
    ):
        selected_guidelines.append(content)

# Create a checkbox for each guideline section
if guidelines:
    with st.container(border=True):
        # Create two columns
        col1, col2 = st.columns(2)

        # Split guidelines into two halves
        guideline_items = list(guidelines.items())
        mid_point = len(guideline_items) // 2

        # First column
        with col1:
            for section_name, content in guideline_items[:mid_point]:
                render_guideline_checkbox(section_name, content, "col1")
                # default = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
                # if st.checkbox(section_name, value=default, key=f"col1_{section_name}"):
                #     selected_guidelines.append(content)

        # Second column
        with col2:
            for section_name, content in guideline_items[mid_point:]:
                render_guideline_checkbox(section_name, content, "col2")
                # default = section_name in ["COMMON GRAMMATICAL ERRORS", "WRITING LETTERS"]
                # if st.checkbox(section_name, value=default, key=f"col2_{section_name}"):
                #     selected_guidelines.append(content)
else:
    st.warning("No guidelines available in the local data.")

# Join all selected guidelines with newlines and store in session state
st.session_state.guidelines = "\n".join(selected_guidelines)

# Show the combined guidelines in a text area
# st.text_area(":blue[**Relevant Guidelines:**]", st.session_state.guidelines, height=200)


# --- NEW helpers: make DOCX/PDF from text ---

def make_docx_bytes(text: str, title: str | None = None) -> bytes:
    """Return a .docx file (bytes) following UKB 04 format structure."""
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
    import re as regex_module
    
    doc = Document()
    
    # Set default font and margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.25)
        section.right_margin = Inches(1.25)
    
    # Helper function to set cell shading
    def set_cell_shading(cell, color_hex):
        """Set cell background color"""
        shading_elm = OxmlElement('w:shd')
        shading_elm.set(qn('w:fill'), color_hex)
        cell._element.get_or_add_tcPr().append(shading_elm)
    
    # Helper function to set cell borders
    def set_cell_border(cell, **kwargs):
        """Set cell borders"""
        tc = cell._element
        tcPr = tc.get_or_add_tcPr()
        tcBorders = OxmlElement('w:tcBorders')
        for edge in ('top', 'left', 'bottom', 'right'):
            if edge in kwargs:
                edge_elem = OxmlElement(f'w:{edge}')
                edge_elem.set(qn('w:val'), 'single')
                edge_elem.set(qn('w:sz'), '4')
                edge_elem.set(qn('w:color'), kwargs[edge])
                tcBorders.append(edge_elem)
        tcPr.append(tcBorders)
    
    # === COVER PAGE ===
    # Add BSP Logo if available
    logo_path = "img/bsp-logo.png"
    if os.path.exists(logo_path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(logo_path, width=Inches(1.5))
    
    doc.add_paragraph()
    
    # BANGKO SENTRAL NG PILIPINAS
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("BANGKO SENTRAL NG PILIPINAS")
    run.font.size = Pt(16)
    run.font.bold = True
    
    doc.add_paragraph()
    
    # FINANCIAL SUPERVISION SECTOR
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("FINANCIAL SUPERVISION SECTOR")
    run.font.size = Pt(12)
    run.font.bold = True
    
    # Add spacing
    for _ in range(3):
        doc.add_paragraph()
    
    # Main Title with box
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("REPORT OF EXAMINATION")
    run.font.size = Pt(16)
    run.font.bold = True
    
    # Add spacing
    for _ in range(2):
        doc.add_paragraph()
    
    # Style/Document Name
    if title:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(title)
        run.font.size = Pt(14)
        run.font.bold = True
    
    doc.add_paragraph()
    
    # Location
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Philippines")
    run.font.size = Pt(11)
    
    # Add spacing
    for _ in range(2):
        doc.add_paragraph()
    
    # Document Type
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Style Rewrite")
    run.font.size = Pt(11)
    
    # Add spacing
    for _ in range(3):
        doc.add_paragraph()
    
    # Date
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    date_str = datetime.now().strftime("%d %B %Y")
    run = p.add_run(f"Date Generated: {date_str}")
    run.font.size = Pt(11)
    
    # Page break
    doc.add_page_break()
    
    # === CONFIDENTIALITY NOTICE PAGE ===
    # Add logo and header
    if os.path.exists(logo_path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(logo_path, width=Inches(1.2))
    
    doc.add_paragraph()
    
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("FINANCIAL SUPERVISION SECTOR")
    run.font.size = Pt(10)
    run.font.bold = True
    
    for _ in range(2):
        doc.add_paragraph()
    
    # Report Title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("REPORT OF EXAMINATION")
    run.font.size = Pt(14)
    run.font.bold = True
    
    doc.add_paragraph()
    
    # Confidentiality Notice
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("THIS REPORT IS STRICTLY CONFIDENTIAL")
    run.font.size = Pt(12)
    run.font.bold = True
    
    doc.add_paragraph()
    
    notice = doc.add_paragraph()
    notice.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    notice_text = ("This report was generated by Bangko Sentral ng Pilipinas (BSP) Style Writer application. "
                   "The content has been rewritten according to the selected editorial style guidelines. "
                   "This document is provided for internal use and review purposes. "
                   "Under no circumstance should this document or any portion thereof be disclosed or made public in any manner, "
                   "except when allowed by law, regulations, or judicial orders. "
                   "Please verify the content for accuracy and compliance before official distribution.")
    run = notice.add_run(notice_text)
    run.font.size = Pt(10)
    
    # Page break
    doc.add_page_break()
    
    # === TABLE OF CONTENTS ===
    doc.add_heading("TABLE OF CONTENTS", level=1)
    doc.add_paragraph()
    
    # TOC entries
    toc_items = [
        ("List of Acronyms", "i"),
        ("Content", "1"),
    ]
    
    for item, page in toc_items:
        p = doc.add_paragraph()
        p.add_run(item).font.size = Pt(11)
        p.add_run("\t").font.size = Pt(11)
        p.add_run(page).font.size = Pt(11)
    
    # Page break
    doc.add_page_break()
    
    # === LIST OF ACRONYMS ===
    doc.add_heading("LIST OF ACRONYMS", level=1)
    doc.add_paragraph()
    
    # Comprehensive list of BSP acronyms
    acronyms = {
        "AC": "Audit Committee",
        "ALCO": "Asset and Liability Committee",
        "ALM": "Asset-Liability Management",
        "AML": "Anti-Money Laundering",
        "AP": "Associated Person",
        "ARA": "Actual Risk Assessment",
        "B2C": "Business-to-Consumer",
        "BAU": "Business-as-Usual",
        "BBS": "Branch Banking Services",
        "BOD": "Board of Directors",
        "BSP": "Bangko Sentral ng Pilipinas",
        "BT": "Bancassurance",
        "CAMEL": "Capital, Assets, Management, Earnings, Liquidity",
        "CASA": "Current and Savings Account",
        "CBS": "Core Banking System",
        "CDD": "Customer Due Diligence",
        "CEO": "Chief Executive Officer",
        "CET": "Common Equity Tier",
        "CFO": "Chief Financial Officer",
        "CIMFS": "Customer Incident Management and Feedback System",
        "CLO": "Chief Lending Officer",
        "CMDI": "Capital Market Development Initiatives",
        "COPC": "Certified Unit Selling Personnel",
        "CORACTS": "Guidelines on Transaction Reporting and Compliance",
        "CRO": "Chief Risk Officer",
        "CTF": "Counter-Terrorism Financing",
        "DCF": "Discounted Cash Flow",
        "DOT": "Declaration of Trust",
        "DST": "Documentary Stamp Tax",
        "EaR": "Earnings at Risk",
        "ECAI": "External Credit Assessment Institution",
        "ECL": "Expected Credit Loss",
        "ECOMM": "E-Commerce",
        "ERM": "Enterprise Risk Management",
        "FMS": "Financial Markets Sector",
        "FOE": "Foreign-Owned Entity",
        "FSS": "Financial Supervision Sector",
        "FVOCI": "Fair Value through Other Comprehensive Income",
        "FVPL": "Fair Value through Profit or Loss",
        "GCG": "Good Corporate Governance",
        "HO": "Head Office",
        "HRMG": "Human Resource Management Group",
        "IAS": "International Accounting Standards",
        "IAASB": "Internal Audit and Regulatory Assessment Process",
        "ICAAP": "Internal Capital Adequacy Assessment Process",
        "IFRS": "International Financial Reporting Standards",
        "IMA": "Investment Management Account",
        "IRRBB": "Interest Rate Risk in the Banking Book",
        "KRI": "Key Risk Indicator",
        "LCR": "Liquidity Coverage Ratio",
        "LGD": "Loss Given Default",
        "LTV": "Loan-to-Value",
        "MIS": "Management Information System",
        "MORB": "Manual of Regulations for Banks",
        "MORNBFI": "Manual of Regulations for Non-Bank Financial Institutions",
        "NII": "Net Interest Income",
        "NIM": "Net Interest Margin",
        "NPL": "Non-Performing Loan",
        "NSFR": "Net Stable Funding Ratio",
        "ORM": "Operational Risk Management",
        "PD": "Probability of Default",
        "PFRS": "Philippine Financial Reporting Standards",
        "RA": "Risk Assessment",
        "RCSA": "Risk and Control Self-Assessment",
        "ROA": "Return on Assets",
        "ROE": "Return on Equity",
        "RP": "Risk Profile",
        "RPT": "Related Party Transaction",
        "RWA": "Risk-Weighted Assets",
        "SME": "Small and Medium Enterprise",
        "TBA": "Treasury Bills Auction",
        "VaR": "Value at Risk",
        "BSFI": "BSP-Supervised Financial Institution",
    }
    
    # Create table for acronyms
    table = doc.add_table(rows=len(acronyms) + 1, cols=2)
    table.style = 'Light Grid Accent 1'
    
    # Header row
    header_cells = table.rows[0].cells
    header_cells[0].text = "Acronym"
    header_cells[1].text = "Definition"
    for cell in header_cells:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(11)
    
    # Data rows
    for idx, (acronym, definition) in enumerate(acronyms.items(), 1):
        row_cells = table.rows[idx].cells
        row_cells[0].text = acronym
        row_cells[1].text = definition
        for cell in row_cells:
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    # Page break
    doc.add_page_break()
    
    # === MAIN CONTENT ===
    doc.add_heading("CONTENT", level=1)
    doc.add_paragraph()
    
    # Helper function to detect table-like content
    def is_table_content(lines):
        """Detect if lines represent a table structure"""
        if len(lines) < 2:
            return False
        non_empty = [line for line in lines if line.strip()]
        if len(non_empty) < 2:
            return False
        # Check for multiple columns indicated by tabs or multiple spaces
        tab_counts = [line.count('\t') for line in non_empty]
        space_pattern = [len(regex_module.findall(r'\s{2,}', line)) for line in non_empty]
        # At least 2 lines must have consistent column separators
        has_tabs = sum(1 for c in tab_counts if c > 0) >= 2
        has_spaces = sum(1 for s in space_pattern if s >= 2) >= 2
        return has_tabs or has_spaces
    
    # Helper function to create formatted table
    def create_assessment_table(lines):
        """Create a formatted BSP assessment table from lines"""
        # Parse table structure
        table_data = []
        for line in lines:
            if line.strip():
                # Split by tab or multiple spaces (2 or more)
                if '\t' in line:
                    cells = [cell.strip() for cell in line.split('\t') if cell.strip()]
                else:
                    # Split by 2+ spaces, filter empty cells
                    cells = [cell.strip() for cell in regex_module.split(r'\s{2,}', line) if cell.strip()]
                if cells:
                    table_data.append(cells)
        
        if not table_data or len(table_data) < 2:
            return None
        
        # Determine number of columns (use most common column count)
        col_counts = [len(row) for row in table_data]
        max_cols = max(col_counts)
        if max_cols == 0 or max_cols == 1:
            return None
        
        # Normalize rows to have consistent column count
        for row in table_data:
            while len(row) < max_cols:
                row.append('')
        
        # Create table with proper styling
        table = doc.add_table(rows=len(table_data), cols=max_cols)
        table.style = 'Light Grid Accent 1'
        table.alignment = WD_TABLE_ALIGNMENT.LEFT
        
        # Set column widths based on content
        if max_cols == 2:
            table.columns[0].width = Inches(4.0)
            table.columns[1].width = Inches(1.5)
        elif max_cols == 3:
            table.columns[0].width = Inches(2.5)
            table.columns[1].width = Inches(2.0)
            table.columns[2].width = Inches(1.5)
        elif max_cols == 4:
            table.columns[0].width = Inches(2.0)
            table.columns[1].width = Inches(2.0)
            table.columns[2].width = Inches(1.0)
            table.columns[3].width = Inches(1.5)
        elif max_cols >= 5:
            for col_idx in range(max_cols):
                table.columns[col_idx].width = Inches(6.0 / max_cols)
        
        # Fill table with data and formatting
        for i, row_data in enumerate(table_data):
            row_cells = table.rows[i].cells
            for j, cell_text in enumerate(row_data):
                cell = row_cells[j]
                # Clear default paragraph
                cell.text = ''
                p = cell.paragraphs[0]
                run = p.add_run(cell_text)
                
                cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                
                # Format first row as header
                if i == 0:
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run.font.bold = True
                    run.font.size = Pt(10)
                    run.font.name = 'Calibri'
                    # Add gray shading to header
                    set_cell_shading(cell, 'D9D9D9')
                else:
                    # Regular cell formatting
                    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    run.font.size = Pt(10)
                    run.font.name = 'Calibri'
                    
                    # Highlight ratings/assessment values
                    if any(keyword in cell_text.upper() for keyword in ['STRONG', 'MODERATE', 'LOW', 'ACCEPTABLE', 'WEAK', 'HIGH']):
                        run.font.bold = True
                
                # Remove extra spacing in paragraphs
                p.space_before = Pt(0)
                p.space_after = Pt(0)
        
        return table
    
    # Process content blocks with enhanced formatting
    current_section = None
    for block in text.replace("\r\n", "\n").split("\n\n"):
        if block.strip():
            stripped_block = block.strip()
            lines = stripped_block.split('\n')
            
            # Check if this is a table structure
            if is_table_content(lines):
                table = create_assessment_table(lines)
                if table:
                    continue
            
            # Detect major section headers
            is_major_header = False
            is_minor_header = False
            
            # Major headers: Roman numerals or risk assessment titles
            if (regex_module.match(r'^(I{1,3}V?|IV|V|VI{0,3}|IX|X{1,3}|XL|L|LX{0,3}|XC|C{1,3})\.\s+', stripped_block) or
                regex_module.match(r'^(Assessment|Directives|Overall|Summary|Scope|Conclusion):', stripped_block, regex_module.IGNORECASE)):
                is_major_header = True
                current_section = stripped_block
            
            # Minor headers: Numbers or labeled items
            elif (regex_module.match(r'^\d+[\.\)]\s+', stripped_block) or
                  (len(lines[0]) < 100 and lines[0].isupper())):
                is_minor_header = True
            
            # Create paragraph with appropriate styling
            if is_major_header:
                # Major section header - larger, bold, with spacing
                p = doc.add_paragraph()
                p.space_before = Pt(12)
                p.space_after = Pt(6)
                run = p.add_run(stripped_block)
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.name = 'Calibri'
            elif is_minor_header:
                # Minor section header - bold, normal size
                p = doc.add_paragraph()
                p.space_before = Pt(6)
                p.space_after = Pt(3)
                run = p.add_run(stripped_block)
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.name = 'Calibri'
            else:
                # Regular content paragraph
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                
                # Handle multi-line content
                for i, line in enumerate(lines):
                    if line.strip():
                        # Check if line itself is a header
                        line_is_header = (len(line.strip()) < 100 and line.strip().isupper())
                        
                        run = p.add_run(line.strip())
                        run.font.size = Pt(11)
                        run.font.name = 'Calibri'
                        
                        if line_is_header:
                            run.font.bold = True
                        
                        # Add line break if not last line
                        if i < len(lines) - 1:
                            p.add_run('\n')
    
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()


def _register_pdf_font_if_available():
    """Optionally register DejaVuSans for better Unicode PDF rendering."""
    try:
        font_path = os.path.join("assets", "DejaVuSans.ttf")
        if os.path.exists(font_path):
            pdfmetrics.registerFont(TTFont("DejaVuSans", font_path))
            return "DejaVuSans"
    except Exception:
        pass
    # Fallback to built-in Helvetica (ASCII/Latin-1 safe)
    return "Helvetica"


def make_pdf_bytes(text: str, title: str | None = None) -> bytes:
    """Return a PDF (bytes) following UKB 04 format structure using ReportLab."""
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
    from reportlab.platypus import PageBreak, Table, TableStyle, Image
    from reportlab.lib import colors
    import re as regex_module
    
    font_name = _register_pdf_font_if_available()

    # Custom page template with header and logo
    def add_page_header(canvas, doc):
        """Add BSP logo and header to each page except cover"""
        canvas.saveState()
        logo_path = "img/bsp-logo.png"
        if os.path.exists(logo_path) and doc.page > 1:
            # Add small logo at top
            canvas.drawImage(logo_path, 2.5*cm, A4[1] - 1.5*cm, width=1.5*cm, height=1.5*cm, preserveAspectRatio=True, mask='auto')
            # Add text next to logo
            canvas.setFont(font_name, 8)
            canvas.drawString(4.5*cm, A4[1] - 1.2*cm, "BANGKO SENTRAL NG PILIPINAS")
            canvas.drawString(4.5*cm, A4[1] - 1.5*cm, "Financial Supervision Sector")
            # Add line
            canvas.setStrokeColor(colors.grey)
            canvas.setLineWidth(0.5)
            canvas.line(2*cm, A4[1] - 2*cm, A4[0] - 2*cm, A4[1] - 2*cm)
        canvas.restoreState()

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=3.2 * cm,
        rightMargin=3.2 * cm,
        topMargin=3 * cm,
        bottomMargin=2.5 * cm,
        title=title or "Rewritten Content",
        author="BSP Style Writer",
    )

    styles = getSampleStyleSheet()
    
    # Custom styles matching UKB format
    bsp_header = ParagraphStyle(
        "BSPHeader",
        fontName=font_name,
        fontSize=16,
        leading=20,
        alignment=TA_CENTER,
        spaceAfter=6,
    )
    
    cover_title = ParagraphStyle(
        "CoverTitle",
        fontName=font_name,
        fontSize=12,
        leading=16,
        alignment=TA_CENTER,
        spaceAfter=18,
        spaceBefore=12,
    )
    
    cover_main = ParagraphStyle(
        "CoverMain",
        fontName=font_name,
        fontSize=16,
        leading=20,
        alignment=TA_CENTER,
        spaceAfter=18,
        spaceBefore=24,
    )
    
    cover_subtitle = ParagraphStyle(
        "CoverSubtitle",
        fontName=font_name,
        fontSize=12,
        leading=16,
        alignment=TA_CENTER,
        spaceAfter=12,
    )
    
    cover_small = ParagraphStyle(
        "CoverSmall",
        fontName=font_name,
        fontSize=11,
        leading=14,
        alignment=TA_CENTER,
        spaceAfter=8,
    )
    
    notice_title = ParagraphStyle(
        "NoticeTitle",
        fontName=font_name,
        fontSize=12,
        leading=16,
        alignment=TA_CENTER,
        spaceAfter=18,
        spaceBefore=6,
    )
    
    notice_body = ParagraphStyle(
        "NoticeBody",
        fontName=font_name,
        fontSize=10,
        leading=14,
        alignment=TA_JUSTIFY,
        spaceAfter=12,
    )
    
    heading_style = ParagraphStyle(
        "CustomHeading",
        fontName=font_name,
        fontSize=14,
        leading=18,
        alignment=TA_LEFT,
        spaceAfter=16,
        spaceBefore=12,
    )
    
    body_style = ParagraphStyle(
        "Body",
        fontName=font_name,
        fontSize=11,
        leading=16,
        alignment=TA_JUSTIFY,
        spaceAfter=12,
    )
    
    toc_style = ParagraphStyle(
        "TOC",
        fontName=font_name,
        fontSize=11,
        leading=16,
        alignment=TA_LEFT,
        spaceAfter=8,
    )

    story = []
    
    # === COVER PAGE ===
    # Add BSP Logo
    logo_path = "img/bsp-logo.png"
    if os.path.exists(logo_path):
        img = Image(logo_path, width=3*cm, height=3*cm)
        img.hAlign = 'CENTER'
        story.append(Spacer(1, 1.5 * cm))
        story.append(img)
        story.append(Spacer(1, 0.5 * cm))
    else:
        story.append(Spacer(1, 2 * cm))
    
    # BSP Header
    story.append(Paragraph("<b>BANGKO SENTRAL NG PILIPINAS</b>", bsp_header))
    story.append(Spacer(1, 0.3 * cm))
    story.append(Paragraph("<b>FINANCIAL SUPERVISION SECTOR</b>", cover_title))
    story.append(Spacer(1, 2.5 * cm))
    
    # Main Title with box effect
    story.append(Paragraph("<b>REPORT OF EXAMINATION</b>", cover_main))
    story.append(Spacer(1, 1.5 * cm))
    
    if title:
        title_safe = title.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        story.append(Paragraph(f"<b>{title_safe}</b>", cover_subtitle))
    
    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph("Philippines", cover_small))
    story.append(Spacer(1, 1 * cm))
    story.append(Paragraph("Style Rewrite", cover_small))
    story.append(Spacer(1, 2 * cm))
    
    date_str = datetime.now().strftime("%d %B %Y")
    story.append(Paragraph(f"Date Generated: {date_str}", cover_small))
    
    # Page break
    story.append(PageBreak())
    
    # === CONFIDENTIALITY NOTICE PAGE ===
    story.append(Spacer(1, 2 * cm))
    story.append(Paragraph("<b>REPORT OF EXAMINATION</b>", heading_style))
    story.append(Spacer(1, 0.5 * cm))
    story.append(Paragraph("<b>THIS REPORT IS STRICTLY CONFIDENTIAL</b>", notice_title))
    story.append(Spacer(1, 0.8 * cm))
    
    notice_text = (
        "This report was generated by Bangko Sentral ng Pilipinas (BSP) Style Writer application. "
        "The content has been rewritten according to the selected editorial style guidelines. "
        "This document is provided for internal use and review purposes. "
        "Under no circumstance should this document or any portion thereof be disclosed or made public in any manner, "
        "except when allowed by law, regulations, or judicial orders. "
        "Please verify the content for accuracy and compliance before official distribution."
    )
    story.append(Paragraph(notice_text, notice_body))
    
    # Page break
    story.append(PageBreak())
    
    # === TABLE OF CONTENTS ===
    story.append(Paragraph("<b>TABLE OF CONTENTS</b>", heading_style))
    story.append(Spacer(1, 0.8 * cm))
    
    toc_data = [
        ["", "Page No."],
        ["List of Acronyms", "i"],
        ["Content", "1"],
    ]
    
    toc_table = Table(toc_data, colWidths=[12*cm, 3*cm])
    toc_table.setStyle(TableStyle([
        ('FONT', (0, 0), (-1, -1), font_name, 11),
        ('FONT', (0, 0), (-1, 0), font_name, 11),
        ('FONTNAME', (0, 0), (-1, 0), font_name),
        ('ALIGN', (1, 0), (1, -1), 'RIGHT'),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
        ('TOPPADDING', (0, 0), (-1, -1), 8),
    ]))
    story.append(toc_table)
    
    # Page break
    story.append(PageBreak())
    
    # === LIST OF ACRONYMS ===
    story.append(Paragraph("<b>LIST OF ACRONYMS</b>", heading_style))
    story.append(Spacer(1, 0.8 * cm))
    
    # Comprehensive list of BSP acronyms
    acronyms_data = [
        ["Acronym", "Definition"],
        ["AC", "Audit Committee"],
        ["ALCO", "Asset and Liability Committee"],
        ["ALM", "Asset-Liability Management"],
        ["AML", "Anti-Money Laundering"],
        ["AP", "Associated Person"],
        ["ARA", "Actual Risk Assessment"],
        ["B2C", "Business-to-Consumer"],
        ["BAU", "Business-as-Usual"],
        ["BBS", "Branch Banking Services"],
        ["BOD", "Board of Directors"],
        ["BSP", "Bangko Sentral ng Pilipinas"],
        ["BT", "Bancassurance"],
        ["CAMEL", "Capital, Assets, Management, Earnings, Liquidity"],
        ["CASA", "Current and Savings Account"],
        ["CBS", "Core Banking System"],
        ["CDD", "Customer Due Diligence"],
        ["CEO", "Chief Executive Officer"],
        ["CET", "Common Equity Tier"],
        ["CFO", "Chief Financial Officer"],
        ["CIMFS", "Customer Incident Management and Feedback System"],
        ["CLO", "Chief Lending Officer"],
        ["CMDI", "Capital Market Development Initiatives"],
        ["COPC", "Certified Unit Selling Personnel"],
        ["CORACTS", "Guidelines on Transaction Reporting and Compliance"],
        ["CRO", "Chief Risk Officer"],
        ["CTF", "Counter-Terrorism Financing"],
        ["DCF", "Discounted Cash Flow"],
        ["DOT", "Declaration of Trust"],
        ["DST", "Documentary Stamp Tax"],
        ["EaR", "Earnings at Risk"],
        ["ECAI", "External Credit Assessment Institution"],
        ["ECL", "Expected Credit Loss"],
        ["ECOMM", "E-Commerce"],
        ["ERM", "Enterprise Risk Management"],
        ["FMS", "Financial Markets Sector"],
        ["FOE", "Foreign-Owned Entity"],
        ["FSS", "Financial Supervision Sector"],
        ["FVOCI", "Fair Value through Other Comprehensive Income"],
        ["FVPL", "Fair Value through Profit or Loss"],
        ["GCG", "Good Corporate Governance"],
        ["HO", "Head Office"],
        ["HRMG", "Human Resource Management Group"],
        ["IAS", "International Accounting Standards"],
        ["IAASB", "Internal Audit and Regulatory Assessment Process"],
        ["ICAAP", "Internal Capital Adequacy Assessment Process"],
        ["IFRS", "International Financial Reporting Standards"],
        ["IMA", "Investment Management Account"],
        ["IRRBB", "Interest Rate Risk in the Banking Book"],
        ["KRI", "Key Risk Indicator"],
        ["LCR", "Liquidity Coverage Ratio"],
        ["LGD", "Loss Given Default"],
        ["LTV", "Loan-to-Value"],
        ["MIS", "Management Information System"],
        ["MORB", "Manual of Regulations for Banks"],
        ["MORNBFI", "Manual of Regulations for Non-Bank Financial Institutions"],
        ["NII", "Net Interest Income"],
        ["NIM", "Net Interest Margin"],
        ["NPL", "Non-Performing Loan"],
        ["NSFR", "Net Stable Funding Ratio"],
        ["ORM", "Operational Risk Management"],
        ["PD", "Probability of Default"],
        ["PFRS", "Philippine Financial Reporting Standards"],
        ["RA", "Risk Assessment"],
        ["RCSA", "Risk and Control Self-Assessment"],
        ["ROA", "Return on Assets"],
        ["ROE", "Return on Equity"],
        ["RP", "Risk Profile"],
        ["RPT", "Related Party Transaction"],
        ["RWA", "Risk-Weighted Assets"],
        ["SME", "Small and Medium Enterprise"],
        ["TBA", "Treasury Bills Auction"],
        ["VaR", "Value at Risk"],
        ["BSFI", "BSP-Supervised Financial Institution"],
    ]
    
    acronyms_table = Table(acronyms_data, colWidths=[3*cm, 12*cm])
    acronyms_table.setStyle(TableStyle([
        ('FONT', (0, 0), (-1, -1), font_name, 9),
        ('FONTNAME', (0, 0), (-1, 0), font_name),
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(acronyms_table)
    
    # Page break
    story.append(PageBreak())
    
    # === MAIN CONTENT ===
    story.append(Paragraph("<b>CONTENT</b>", heading_style))
    story.append(Spacer(1, 0.5 * cm))
    
    # Header detection style
    header_style = ParagraphStyle(
        "HeaderText",
        fontName=font_name,
        fontSize=11,
        leading=16,
        alignment=TA_JUSTIFY,
        spaceAfter=12,
        spaceBefore=6,
    )
    
    # Process content blocks with header detection
    for block in text.replace("\r\n", "\n").split("\n\n"):
        if block.strip():
            stripped_block = block.strip()
            
            # Detect headers: Roman numerals, numbers, or ALL CAPS lines
            is_header = False
            if (regex_module.match(r'^(I{1,3}V?|IV|V|VI{0,3}|IX|X{1,3}|XL|L|LX{0,3}|XC|C{1,3})\.\s+', stripped_block) or
                regex_module.match(r'^\d+[\.\)]\s+', stripped_block) or
                (len(stripped_block.split('\n')[0]) < 100 and stripped_block.split('\n')[0].isupper())):
                is_header = True
            
            # Escape HTML-sensitive characters but preserve structure
            block_safe = stripped_block.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            
            # Process line by line to handle mixed content
            lines = block_safe.split('\n')
            formatted_lines = []
            for line in lines:
                if line.strip():
                    # Check if individual line is a header
                    line_is_header = (len(line.strip()) < 100 and line.strip().isupper()) or is_header
                    if line_is_header:
                        formatted_lines.append(f"<b>{line.strip()}</b>")
                    else:
                        formatted_lines.append(line.strip())
            
            final_text = "<br/>".join(formatted_lines)
            
            # Use appropriate style based on content type
            if is_header:
                story.append(Paragraph(final_text, header_style))
            else:
                story.append(Paragraph(final_text, body_style))

    doc.build(story, onFirstPage=add_page_header, onLaterPages=add_page_header)
    return buf.getvalue()


if st.button(
    ":blue[**Rewrite Content**]",
    key="extract",
    disabled=(
        content_all.strip() == ""
        or st.session_state.style == ""
        or st.session_state.example == ""
    ),
):
    with st.spinner("Processing..."):
        # --- Process and store the result ---
        output = prompts.rewrite_content(content_all, max_output_length, False)
        utils.save_output(output, content_all)

        # --- Store in session state ---
        st.session_state["last_output"] = output
        ts = datetime.now().strftime("%Y%m%d-%H%M%S")
        style_id = (st.session_state.get("styleId") or "Style").replace(" ", "_")
        base_name = f"rewrite_{style_id}_{ts}"
        st.session_state["last_base_name"] = base_name
        st.session_state["last_title"] = f"Rewrite • {st.session_state.get('styleId') or 'Selected Style'}"
        st.session_state["output_ready"] = True
        st.rerun()

# Display output and download buttons if available (only once)
if st.session_state.get("output_ready") and st.session_state.get("last_output"):
    with st.container(border=True):
        st.markdown("### ✨ Rewritten Output")
        
        # Show the output text
        st.text_area(
            "Result",
            st.session_state["last_output"],
            height=300,
            key="output_display",
        )
        
        # Generate download files
        title_text = st.session_state.get("last_title", "Rewrite")
        base_name = st.session_state.get("last_base_name", "rewrite")
        
        docx_bytes = make_docx_bytes(st.session_state["last_output"], title=title_text)
        pdf_bytes = make_pdf_bytes(st.session_state["last_output"], title=title_text)

        c1, c2 = st.columns(2)
        with c1:
            st.download_button(
                "⬇️ Download as DOCX",
                data=docx_bytes,
                file_name=f"{base_name}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                width='stretch',
                key="download_docx",
            )
        with c2:
            st.download_button(
                "⬇️ Download as PDF",
                data=pdf_bytes,
                file_name=f"{base_name}.pdf",
                mime="application/pdf",
                width='stretch',
                key="download_pdf",
            )